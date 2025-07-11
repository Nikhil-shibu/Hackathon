"""
Script to generate a PowerPoint presentation for the Thrive Mind project.
Requires: python-pptx (install via `pip install python-pptx`)

Run: python generate_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Thrive Mind"
slide.placeholders[1].text = (
    "A supportive companion for individuals with Autism, Down Syndrome, Dementia, "
    "and learning disabilities."
)

# Function to add content slide
def add_content_slide(title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            tf.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 1

# Add slides
add_content_slide("Overview", [
    "Purpose: Provide routines, communication tools, and safety features.",
    "Target Users: Individuals with cognitive or developmental challenges.",
    "Platforms: Flutter (Web, Mobile, Desktop)."
])

add_content_slide("Key Features", [
    "Routine Management: Schedule and reminders for daily tasks.",
    "Communication Aid: Text-to-speech and speech-to-text integration.",
    "Safety Alerts: Emergency notifications and location services.",
    "Accessibility: Large text, high contrast, voice guidance."
])

add_content_slide("Architecture", [
    "Flutter Frontend: Single codebase for multiple platforms.",
    "Firebase Backend: Authentication, Firestore, Notifications.",
    "State Management: Provider package."
])

add_content_slide("Authentication Flow", [
    "Login: Email/password with Firebase Auth.",
    "Persistent Login: Session maintained until logout.",
    "Logout: Clears session and navigates to login screen."
])

add_content_slide("ML Integration Plan", [
    "Use Case: Predict when a user might miss a routine.",
    "Data: Historical completion timestamps, delays, user context.",
    "Model: Time-series classifier, exported to TensorFlow Lite.",
    "Integration: Inference via tflite_flutter plugin."
])

add_content_slide("Next Steps", [
    "Enhance ML accuracy with more data and retraining cycles.",
    "Add user customization and preferences.",
    "Perform usability testing and collect feedback.",
    "Deploy to app stores and monitor analytics."
])

# Save presentation
prs.save('ThriveMind_Presentation.pptx')
Requires: python-pptx (install via `pip install python-pptx`)

Run: python generate_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Helper to add slides
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    try:
        slide.placeholders[1].text = subtitle
    except Exception:
        # fallback: add textbox
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = subtitle

# Helper to add title/content slides
from pptx.enum.text import PP_ALIGN

def add_content_slide(title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    try:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for idx, line in enumerate(content_lines):
            if idx == 0:
                p = tf
                p.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 1
    except Exception:
        # fallback: full slide textbox
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = txBox.text_frame
        for idx, line in enumerate(content_lines):
            p = tf.add_paragraph() if idx else tf
            p.text = line

# Generate slides
add_title_slide(
    "Thrive Mind",
    "A supportive companion for individuals with Autism, Down Syndrome, Dementia, and learning disabilities."
)

add_content_slide("Overview", [
    "Purpose: Provide routines, communication tools, and safety features.",
    "Target Users: Individuals with cognitive or developmental challenges.",
    "Platforms: Flutter (Web, Mobile, Desktop)."
])

add_content_slide("Key Features", [
    "Routine Management: Schedule and reminders for daily tasks.",
    "Communication Aid: Text-to-speech and speech-to-text integration.",
    "Safety Alerts: Emergency notifications and location services.",
    "Accessibility: Large text, high contrast, voice guidance."
])

add_content_slide("Architecture", [
    "Flutter Frontend: Single codebase for multiple platforms.",
    "Firebase Backend: Authentication, Firestore, Notifications.",
    "Providers: State management via Provider package."
])

add_content_slide("Authentication Flow", [
    "Login: Email/password with Firebase Auth.",
    "Persistent Login: Session maintained until logout.",
    "Logout: Clears session and navigates to login screen."
])

add_content_slide("ML Integration Plan", [
    "Use Case: Predict when a user might miss a routine.",
    "Data: Historical completion timestamps, delays, user context.",
    "Model: Train time-series classifier, export to TensorFlow Lite.",
    "Integration: Inference via tflite_flutter plugin."
])

add_content_slide("Next Steps", [
    "Enhance ML accuracy with more data and retraining cycles.",
    "Add user customization and preferences.",
    "Perform usability testing and collect feedback.",
    "Deploy to app stores and monitor analytics."
])

# Save the presentation
prs.save('ThriveMind_Presentation.pptx')
prs.save('ThriveMind_Presentation.pptx')

slide.shapes.title.text = "Thrive Mind"
slide.placeholders[1].text = "A supportive companion for individuals with Autism, Down Syndrome, Dementia, and learning disabilities."

# Slide 2: Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Overview"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Purpose: Provide routines, communication tools, and safety features."
p = tf.add_paragraph()
p.text = "Target Users: Individuals with cognitive or developmental challenges."
p.level = 1
p = tf.add_paragraph()
p.text = "Platforms: Flutter (Web, Mobile, Desktop)."
p.level = 1

# Slide 3: Key Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Features"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Routine Management: Schedule and reminders for daily tasks."
p = tf.add_paragraph()
p.text = "Communication Aid: Text-to-speech and speech-to-text integration."
p = tf.add_paragraph()
p.text = "Safety Alerts: Emergency notifications and location services."
p = tf.add_paragraph()
p.text = "Accessibility: Large text, high contrast, voice guidance."

# Slide 4: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Architecture"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Flutter Frontend: Single codebase for multiple platforms."
p = tf.add_paragraph()
p.text = "Firebase Backend: Authentication, Firestore, Notifications."
p = tf.add_paragraph()
p.text = "Providers: State management via Provider package."

# Slide 5: Authentication Flow
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Authentication Flow"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Login: Email/password with Firebase Auth."
p = tf.add_paragraph()
p.text = "Persistent Login: Session maintained until logout."
p = tf.add_paragraph()
p.text = "Logout: Clears session and navigates to login screen."

# Slide 6: ML Integration Plan
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "ML Integration Plan"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Use Case: Predict when a user might miss a routine."
p = tf.add_paragraph()
p.text = "Data: Historical completion timestamps, delays, user context."
p = tf.add_paragraph()
p.text = "Model: Train time-series classifier, export to TensorFlow Lite."
p = tf.add_paragraph()
p.text = "Integration: Inference via tflite_flutter plugin."

# Slide 7: Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Next Steps"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Enhance ML accuracy with more data and retraining cycles."
p = tf.add_paragraph()
p.text = "Add user customization and preferences."
p = tf.add_paragraph()
p.text = "Perform usability testing and collect feedback."
p = tf.add_paragraph()
p.text = "Deploy to app stores and monitor analytics."

# Save the presentation
prs.save('ThriveMind_Presentation.pptx')
